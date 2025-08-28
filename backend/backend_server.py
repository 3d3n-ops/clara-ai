import os
import asyncio
import logging
import redis.asyncio as redis
import httpx  # Ensure this is imported for async HTTP requests

from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse

import PyPDF2
import docx
from pptx import Presentation
import openai
from dotenv import load_dotenv

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('backend_server.log'),  # Log to file
        logging.StreamHandler()  # Also log to console
    ]
)
logger = logging.getLogger(__name__)

# Load environment variables
load_dotenv()

# Initialize FastAPI app
app = FastAPI()

# CORS Configuration
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Allows all origins
    allow_credentials=True,
    allow_methods=["*"],  # Allows all methods
    allow_headers=["*"],  # Allows all headers
)

# Initialize Redis client
try:
    redis_url = os.getenv("REDIS_URL")
    if not redis_url:
        logger.warning("REDIS_URL not set, using default localhost URL")
        redis_url = "redis://localhost:6379"
    
    redis_client = redis.from_url(redis_url, encoding="utf-8", decode_responses=True)
    logger.info(f"Redis client initialized with URL: {redis_url}")
except Exception as e:
    logger.error(f"Failed to initialize Redis client: {e}")
    redis_client = None

# Modify OpenAI client initialization
try:
    openai.api_key = os.getenv("OPENAI_API_KEY")
    if not openai.api_key:
        logger.error("OpenAI API key not found in environment variables")
    
    # Use async client
    openai.async_client = openai.AsyncOpenAI(
        api_key=openai.api_key,
        http_client=httpx.AsyncClient()
    )
except Exception as e:
    logger.error(f"Error setting up OpenAI async client: {e}")
    openai.async_client = None

def extract_text_from_pdf(file) -> str:
    """Extract text from PDF file."""
    try:
        logger.info("Extracting text from PDF file")
        reader = PyPDF2.PdfReader(file)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        logger.info(f"Extracted {len(text)} characters from PDF")
        return text
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {e}")
        raise

def extract_text_from_docx(file) -> str:
    """Extract text from DOCX file."""
    try:
        logger.info("Extracting text from DOCX file")
        doc = docx.Document(file)
        text = "\n".join([p.text for p in doc.paragraphs])
        logger.info(f"Extracted {len(text)} characters from DOCX")
        return text
    except Exception as e:
        logger.error(f"Error extracting text from DOCX: {e}")
        raise

def extract_text_from_ppt(file) -> str:
    """Extract text from PPT file."""
    try:
        logger.info("Extracting text from PPT file")
        presentation = Presentation(file)
        text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text.append(run.text)
        extracted_text = "\n".join(text)
        logger.info(f"Extracted {len(extracted_text)} characters from PPT")
        return extracted_text
    except Exception as e:
        logger.error(f"Error extracting text from PPT: {e}")
        raise

@app.post("/file-upload/")
async def upload_file(file: UploadFile = File(...)):
    """
    Upload and process file, generate notes, diagram, and title.
    Store summary in Redis cache.
    """
    try:
        # Validate file size (10MB limit)
        file.file.seek(0, 2)  # Move to end of file
        file_size = file.file.tell()
        file.file.seek(0)  # Reset file pointer
        
        MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB
        if file_size > MAX_FILE_SIZE:
            raise HTTPException(
                status_code=400, 
                detail=f"File size exceeds maximum limit of 10MB. Current size: {file_size/1024/1024:.2f}MB"
            )
        
        # Validate file extension
        if not file.filename:
            raise HTTPException(status_code=400, detail="Invalid file name")
        
        file_ext = file.filename.lower().split('.')[-1]
        ALLOWED_EXTENSIONS = {'pdf', 'docx', 'pptx', 'txt', 'png', 'jpg', 'jpeg'}
        
        if file_ext not in ALLOWED_EXTENSIONS:
            raise HTTPException(
                status_code=400, 
                detail=f"Unsupported file type. Allowed types: {', '.join(ALLOWED_EXTENSIONS)}"
            )
        
        logger.info(f"Received file upload: {file.filename}")
        
        # Determine file type and extract text
        if file.filename.endswith(".pdf"):
            logger.info("Detected PDF file")
            text = extract_text_from_pdf(file.file)
        elif file.filename.endswith(".docx"):
            logger.info("Detected DOCX file")
            text = extract_text_from_docx(file.file)
        elif file.filename.endswith(".pptx"):
            logger.info("Detected PPTX file")
            text = extract_text_from_ppt(file.file)
        elif file.filename.lower().endswith(('.png', '.jpg', '.jpeg', '.txt')):
            # For images and text files, we'll use a simple text extraction
            text = await file.read()
            text = text.decode('utf-8') if isinstance(text, bytes) else str(text)
        else:
            logger.warning(f"Unsupported file format: {file.filename}")
            raise HTTPException(status_code=400, detail="Unsupported file format")

        logger.info(f"Extracted {len(text)} total characters from file")

        # Async tasks for generating content
        logger.info("Starting OpenAI tasks for notes, diagram, and title generation")
        
        try:
            # Use async_client for API calls
            notes_response = await openai.async_client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {"role": "system", "content": "You are a helpful note-taker."},
                    {"role": "user", "content": f"Make clear, concise study notes on the following text: \n\n{text}"}
                ],
                max_tokens=500
            )

            diagram_response = await openai.async_client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {"role": "system", "content": "You are a diagram generator for visual learners."},
                    {"role": "user", "content": f"Make clear and understandable diagrams on the following text: \n\n{text}"}
                ],
                max_tokens=300
            )

            title_response = await openai.async_client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {"role": "system", "content": "You are a helpful note-title generator."},
                    {"role": "user", "content": f"Make simple title on the following text: \n\n{text}"}
                ],
                max_tokens=50
            )

            # Extract results with additional validation
            notes = notes_response.choices[0].message.content or "No notes generated"
            diagram = diagram_response.choices[0].message.content or "No diagram generated"
            title = title_response.choices[0].message.content or "Untitled"

            logger.info("Successfully generated notes, diagram, and title")

        except Exception as e:
            logger.error(f"Error generating content with OpenAI: {e}", exc_info=True)
            notes = "Error generating notes"
            diagram = "Error generating diagram"
            title = "Error generating title"

        # Store summary in Redis with a unique key
        summary_key = f"summary:{hash(text)}"
        
        if redis_client:
            try:
                await redis_client.hset(summary_key, mapping={
                    "title": title,
                    "notes": notes,
                    "diagram": diagram
                })
                # Optional: Set expiration for the key (e.g., 1 hour)
                await redis_client.expire(summary_key, 3600)
                logger.info(f"Stored summary in Redis with key: {summary_key}")
            except Exception as redis_error:
                logger.error(f"Failed to store summary in Redis: {redis_error}")
        else:
            logger.warning("Redis client not initialized, skipping cache storage")

        return JSONResponse(content={
            "title": title,
            "notes": notes,
            "diagram": diagram,
            "cache_key": summary_key
        })

    except HTTPException as http_error:
        logger.error(f"HTTP Error: {http_error.detail}")
        raise
    except Exception as e:
        # Comprehensive error handling
        logger.error(f"Unexpected error during file upload: {e}", exc_info=True)
        return JSONResponse(
            status_code=500, 
            content={"error": str(e), "message": "An error occurred while processing the file"}
        )

# Optional: Add a route to retrieve cached summary
@app.get("/summary/{cache_key}")
async def get_summary(cache_key: str):
    """Retrieve a previously generated summary from Redis cache."""
    try:
        logger.info(f"Attempting to retrieve summary for key: {cache_key}")
        
        if not redis_client:
            logger.error("Redis client not initialized")
            raise HTTPException(status_code=500, detail="Redis client not available")
        
        summary = await redis_client.hgetall(cache_key)
        
        if not summary:
            logger.warning(f"No summary found for key: {cache_key}")
            raise HTTPException(status_code=404, detail="Summary not found")
        
        logger.info(f"Successfully retrieved summary for key: {cache_key}")
        return summary
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error retrieving summary: {e}", exc_info=True)
        return JSONResponse(
            status_code=500, 
            content={"error": str(e), "message": "An error occurred while retrieving the summary"}
        )

# Optional: Startup event to verify connections
@app.on_event("startup")
async def startup_event():
    """Verify connections on startup."""
    logger.info("Application starting up")
    
    # Verify Redis connection
    if redis_client:
        try:
            await redis_client.ping()
            logger.info("Successfully connected to Redis")
        except Exception as e:
            logger.error(f"Failed to connect to Redis: {e}")
    
    # Verify OpenAI API key
    try:
        # A simple test to verify OpenAI API key
        openai.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": "Test connection"}],
            max_tokens=10
        )
        logger.info("Successfully verified OpenAI API connection")
    except Exception as e:
        logger.error(f"Failed to verify OpenAI API connection: {e}")

# Optional: Shutdown event for cleanup
@app.on_event("shutdown")
async def shutdown_event():
    """Perform cleanup on application shutdown."""
    logger.info("Application shutting down")
    
    if redis_client:
        try:
            await redis_client.close()
            logger.info("Redis client closed successfully")
        except Exception as e:
            logger.error(f"Error closing Redis client: {e}") 
